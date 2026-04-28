"""PPTX export.

Coordinates are already in EMU, which python-pptx uses natively (Emu()).
We add slides, apply backgrounds, sort elements by zIndex, and emit
editable text/image/shape/line elements.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Pt

from .models import PresentationDocument, SlideDocument, SlideElement
from .storage import asset_dir

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

SHAPE_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
}

FILL_TAGS = {
    qn("a:noFill"),
    qn("a:solidFill"),
    qn("a:gradFill"),
    qn("a:blipFill"),
    qn("a:pattFill"),
    qn("a:grpFill"),
}


def _hex_value(value: str | None, fallback: str = "FFFFFF") -> str:
    v = (value or "").strip().lstrip("#")
    if len(v) == 3:
        v = "".join(ch * 2 for ch in v)
    if len(v) != 6:
        return fallback
    try:
        int(v, 16)
    except ValueError:
        return fallback
    return v.upper()


def _hex_to_rgb(value: str) -> RGBColor:
    v = _hex_value(value)
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def _resolve_asset_path(asset_id: str | None, doc: PresentationDocument) -> Path | None:
    if not asset_id:
        return None
    asset = next((a for a in doc.assets if a.id == asset_id), None)
    folder = asset_dir(asset_id)
    if not asset or not folder:
        return None
    if asset.fileName:
        candidate = folder / asset.fileName
        if candidate.exists():
            return candidate
    files = [p for p in folder.iterdir() if p.is_file()]
    return files[0] if files else None


def _apply_background(slide, bg, doc: PresentationDocument) -> None:
    if bg.type == "color":
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(bg.value)
    elif bg.type == "gradient":
        _add_gradient_rect(
            slide,
            doc,
            getattr(bg, "from_", "#FFFFFF"),
            bg.to,
            bg.angle,
        )
    elif bg.type == "image":
        path = _resolve_asset_path(bg.assetId, doc)
        if path is not None:
            slide.shapes.add_picture(
                str(path), 0, 0, width=Emu(doc.slideSize.widthEmu), height=Emu(doc.slideSize.heightEmu)
            )


def _set_shape_gradient_fill(shape, from_color: str, to_color: str, angle: float) -> None:
    sp_pr = shape._element.spPr
    for child in list(sp_pr):
        if child.tag in FILL_TAGS:
            sp_pr.remove(child)

    from_hex = _hex_value(from_color)
    to_hex = _hex_value(to_color, "111827")
    angle_units = int((angle % 360) * 60000)
    grad_fill = parse_xml(
        f"""
        <a:gradFill {nsdecls("a")} flip="none" rotWithShape="1">
          <a:gsLst>
            <a:gs pos="0">
              <a:srgbClr val="{from_hex}"/>
            </a:gs>
            <a:gs pos="100000">
              <a:srgbClr val="{to_hex}"/>
            </a:gs>
          </a:gsLst>
          <a:lin ang="{angle_units}" scaled="1"/>
        </a:gradFill>
        """
    )

    insert_at = len(sp_pr)
    for idx, child in enumerate(sp_pr):
        if child.tag == qn("a:ln"):
            insert_at = idx
            break
    sp_pr.insert(insert_at, grad_fill)


def _add_gradient_rect(
    slide,
    doc: PresentationDocument,
    from_color: str,
    to_color: str,
    angle: float,
) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Emu(doc.slideSize.widthEmu),
        Emu(doc.slideSize.heightEmu),
    )
    _set_shape_gradient_fill(bg, from_color, to_color, angle)
    bg.line.fill.background()


def _add_text(slide, el, theme_text_color: str) -> None:
    box = slide.shapes.add_textbox(
        Emu(el.frame.xEmu),
        Emu(el.frame.yEmu),
        Emu(el.frame.wEmu),
        Emu(el.frame.hEmu),
    )
    tf = box.text_frame
    tf.word_wrap = True
    lines = (el.text or "").split("\n") if el.text else [""]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text_value = line
        if el.role == "bulletList":
            text_value = line.lstrip().removeprefix("•").removeprefix("-").removeprefix("*").strip()
        if el.role == "bulletList" and text_value:
            _enable_bullet(para)
        para.alignment = ALIGN_MAP.get(el.style.align, PP_ALIGN.LEFT)
        run = para.add_run()
        run.text = text_value
        run.font.name = el.style.fontFamily
        run.font.size = Pt(el.style.fontSize)
        if el.style.fontWeight and el.style.fontWeight >= 600:
            run.font.bold = True
        if el.style.italic:
            run.font.italic = True
        if el.style.underline:
            run.font.underline = True
        run.font.color.rgb = _hex_to_rgb(el.style.color or theme_text_color)


def _enable_bullet(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag in {qn("a:buNone"), qn("a:buAutoNum"), qn("a:buChar"), qn("a:buBlip")}:
            p_pr.remove(child)
    p_pr.insert(0, parse_xml(f'<a:buChar {nsdecls("a")} char="•"/>'))


def _add_image(slide, el, doc: PresentationDocument) -> None:
    path = _resolve_asset_path(el.assetId, doc)
    if path is None:
        return
    slide.shapes.add_picture(
        str(path),
        Emu(el.frame.xEmu),
        Emu(el.frame.yEmu),
        width=Emu(el.frame.wEmu),
        height=Emu(el.frame.hEmu),
    )


def _add_shape(slide, el) -> None:
    shape_type = SHAPE_MAP.get(el.shape, MSO_SHAPE.RECTANGLE)
    sh = slide.shapes.add_shape(
        shape_type,
        Emu(el.frame.xEmu),
        Emu(el.frame.yEmu),
        Emu(el.frame.wEmu),
        Emu(el.frame.hEmu),
    )
    if el.style.fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = _hex_to_rgb(el.style.fill)
    else:
        sh.fill.background()
    if el.style.stroke:
        sh.line.color.rgb = _hex_to_rgb(el.style.stroke)
        if el.style.strokeWidth:
            sh.line.width = Emu(el.style.strokeWidth)
    else:
        sh.line.fill.background()


def _add_line(slide, el) -> None:
    connector = slide.shapes.add_connector(
        1,  # straight
        Emu(el.frame.xEmu),
        Emu(el.frame.yEmu),
        Emu(el.frame.xEmu + el.frame.wEmu),
        Emu(el.frame.yEmu + el.frame.hEmu),
    )
    connector.line.color.rgb = _hex_to_rgb(el.style.color)
    connector.line.width = Emu(el.style.widthEmu)


def _emit(slide, el: SlideElement, doc: PresentationDocument) -> None:
    if el.type == "text":
        _add_text(slide, el, doc.theme.colors.text)
    elif el.type == "image":
        _add_image(slide, el, doc)
    elif el.type == "shape":
        _add_shape(slide, el)
    elif el.type == "line":
        _add_line(slide, el)
    elif el.type == "icon":
        # icons fall back to image rendering when an asset is provided
        _add_image(slide, el, doc)
    elif el.type == "group":
        # groups are exploded in v1
        return


def export_to_pptx(doc: PresentationDocument) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(doc.slideSize.widthEmu)
    prs.slide_height = Emu(doc.slideSize.heightEmu)
    blank_layout = prs.slide_layouts[6]  # blank
    for slide_doc in doc.slides:
        slide = prs.slides.add_slide(blank_layout)
        _apply_background(slide, slide_doc.background, doc)
        ordered = sorted(slide_doc.elements, key=lambda e: e.zIndex)
        for el in ordered:
            _emit(slide, el, doc)
        if slide_doc.notes:
            slide.notes_slide.notes_text_frame.text = slide_doc.notes
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
