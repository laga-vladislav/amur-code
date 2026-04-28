from __future__ import annotations

import io
from typing import Any

from fastapi import UploadFile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .schemas import RawPresentation
from .utils import int_value, rgb_to_hex

import tempfile
# def _open_presentation(source: str | UploadFile | bytes | bytearray) -> Presentation:
#     if isinstance(source, UploadFile):
#         file_obj = source.file
#         if hasattr(file_obj, "file"):
#             file_obj = file_obj.file
#         content = file_obj.read()
#         return Presentation(io.BytesIO(content))
#     if isinstance(source, (bytes, bytearray)):
#         return Presentation(io.BytesIO(source))
#     return Presentation(source)


# def _open_presentation(source: str | UploadFile | bytes | bytearray) -> Presentation:
#     if isinstance(source, UploadFile):
#         content = source.file.read()
#         return Presentation(io.BytesIO(content))

#     if isinstance(source, (bytes, bytearray)):
#         return Presentation(io.BytesIO(source))

#     return Presentation(source)


def _open_presentation(source: UploadFile) -> Presentation:
    with tempfile.NamedTemporaryFile(suffix=".pptx") as tmp:
        tmp.write(source.file.read())
        tmp.flush()
        return Presentation(tmp.name)

def parse_pptx(source: str | UploadFile | bytes | bytearray) -> dict[str, Any]:
    prs = _open_presentation(source)
    slides_data: list[dict[str, Any]] = []

    for index, slide in enumerate(prs.slides):
        elements: list[dict[str, Any]] = []

        for shape in slide.shapes:
            bbox = {
                "x": int_value(shape.left),
                "y": int_value(shape.top),
                "width": int_value(shape.width),
                "height": int_value(shape.height),
            }
            element: dict[str, Any] = {
                "type": "unknown",
                "text": None,
                "bbox": bbox,
                "style": {
                    "fontSize": None,
                    "fontFamily": None,
                    "bold": False,
                    "color": None,
                },
                "isBullet": False,
            }

            if getattr(shape, "has_text_frame", False):
                text_frame = shape.text_frame
                element["type"] = "text"
                element["text"] = shape.text or None

                if text_frame.paragraphs:
                    paragraph = text_frame.paragraphs[0]
                    if paragraph.runs:
                        run = paragraph.runs[0]
                        font = run.font
                        font_size = None
                        if font.size is not None:
                            font_size = int(round(font.size.pt))
                        element["style"] = {
                            "fontSize": font_size,
                            "fontFamily": font.name,
                            "bold": bool(font.bold),
                            "color": rgb_to_hex(font.color.rgb) if getattr(font.color, "rgb", None) else None,
                        }

                    # element["isBullet"] = any(
                    #     paragraph.level > 0
                    #     or bool(getattr(paragraph.paragraph_format, "bullet", False))
                    #     for paragraph in text_frame.paragraphs
                    # )

                    element["isBullet"] = any(
                        paragraph.level is not None and paragraph.level > 0
                        for paragraph in text_frame.paragraphs
                    )

                    
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                element["type"] = "image"

            elements.append(element)

        slides_data.append({"index": index, "elements": elements})

    raw = {
        "slideSize": {
            "widthEmu": int_value(prs.slide_width),
            "heightEmu": int_value(prs.slide_height),
        },
        "slides": slides_data,
    }
    RawPresentation.model_validate(raw)
    return raw
