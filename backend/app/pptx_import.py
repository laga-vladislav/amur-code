"""Import a .pptx file as a TemplateDocument.

Each PowerPoint slide becomes a layout. We extract slide size, basic theme
colors (best-effort), and convert shapes into our element types:

* text frames → text elements (with the role inferred from font size /
  position).
* placeholders → placeholder text elements that the AI generator can fill.
* picture shapes → image placeholder elements (we do not extract bitmap
  bytes — the user re-uploads images if needed).
* auto-shapes → shape elements (rect / roundRect / ellipse / triangle).
* connectors → line elements.

The goal is to produce a usable layout the user can refine in the editor,
not a 1:1 fidelity import.
"""

from __future__ import annotations

import io
import uuid
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from .models import (
    BackgroundColor,
    LayoutDocument,
    PresentationTheme,
    SlideSize,
    TemplateDocument,
    ThemeColors,
    ThemeFonts,
)


SLIDE_TYPES_PRIORITY = ("title", "bullets", "image", "text", "conclusion")


def _new_id(prefix: str) -> str:
    return f"{prefix}_{uuid.uuid4().hex[:10]}"


def _hex_from_rgb(rgb) -> str | None:
    try:
        if rgb is None:
            return None
        return f"#{int(rgb):06X}"
    except (TypeError, ValueError):
        return None


def _theme_colors_from_pptx(prs: Presentation) -> ThemeColors:
    colors = ThemeColors(
        background="#FFFFFF",
        text="#111827",
        primary="#2563EB",
        secondary="#64748B",
        accent="#F97316",
    )
    try:
        scheme = prs.slide_master.element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme"
        )
        if scheme is None:
            return colors
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        mapping = {
            "lt1": "background",
            "dk1": "text",
            "accent1": "primary",
            "accent2": "secondary",
            "accent3": "accent",
        }
        updates: dict[str, str] = {}
        for child in scheme:
            tag = child.tag.removeprefix(ns)
            attr = mapping.get(tag)
            if not attr:
                continue
            srgb = child.find(f"{ns}srgbClr")
            sys = child.find(f"{ns}sysClr")
            value = None
            if srgb is not None and srgb.get("val"):
                value = f"#{srgb.get('val').upper()}"
            elif sys is not None and sys.get("lastClr"):
                value = f"#{sys.get('lastClr').upper()}"
            if value:
                updates[attr] = value
        if updates:
            colors = colors.model_copy(update=updates)
    except Exception:
        pass
    return colors


def _slide_background_color(slide, fallback: str) -> BackgroundColor:
    try:
        fill = slide.background.fill
        if fill.type == 1:  # MSO_FILL_TYPE.SOLID
            hx = _hex_from_rgb(fill.fore_color.rgb)
            if hx:
                return BackgroundColor(type="color", value=hx)
    except Exception:
        pass
    return BackgroundColor(type="color", value=fallback)


def _infer_text_role(text_frame, frame_emu: dict[str, int]) -> str:
    """Best-effort role detection based on size and y position."""
    text = (text_frame.text or "").strip()
    sizes: list[int] = []
    has_bullets = False
    paragraph_count = 0
    for para in text_frame.paragraphs:
        if (para.text or "").strip():
            paragraph_count += 1
        if para.level and para.level > 0:
            has_bullets = True
        for run in para.runs:
            if run.font.size:
                try:
                    sizes.append(int(run.font.size.pt))
                except Exception:
                    pass
    max_size = max(sizes) if sizes else 0
    y = frame_emu.get("yEmu", 0)
    h = frame_emu.get("hEmu", 0)

    if max_size >= 32 or (y < 1500000 and h < 2000000 and paragraph_count == 1 and len(text) < 120):
        return "title"
    if has_bullets or paragraph_count >= 3 or text.count("\n") >= 2:
        return "bulletList"
    if max_size >= 22 and paragraph_count == 1:
        return "subtitle"
    return "body"


def _font_from_text_frame(text_frame, default_color: str) -> dict[str, Any]:
    sizes: list[int] = []
    color = default_color
    bold = False
    italic = False
    align = "left"
    family = "Inter"
    for para in text_frame.paragraphs:
        if para.alignment is not None:
            mapping = {1: "left", 2: "center", 3: "right", 4: "justify"}
            align = mapping.get(int(para.alignment), align)
        for run in para.runs:
            if run.font.size:
                try:
                    sizes.append(int(run.font.size.pt))
                except Exception:
                    pass
            if run.font.bold:
                bold = True
            if run.font.italic:
                italic = True
            try:
                rgb = run.font.color.rgb if run.font.color and run.font.color.type else None
                if rgb is not None:
                    hx = _hex_from_rgb(rgb)
                    if hx:
                        color = hx
            except Exception:
                pass
            if run.font.name:
                family = run.font.name
    font_size = max(sizes) if sizes else 18
    return {
        "fontFamily": family,
        "fontSize": font_size,
        "fontWeight": 700 if bold else 400,
        "italic": italic,
        "color": color,
        "align": align,
        "lineHeight": 1.3,
    }


def _frame_from_shape(shape) -> dict[str, int]:
    return {
        "xEmu": int(getattr(shape, "left", 0) or 0),
        "yEmu": int(getattr(shape, "top", 0) or 0),
        "wEmu": int(getattr(shape, "width", 914400) or 914400),
        "hEmu": int(getattr(shape, "height", 914400) or 914400),
        "rotate": float(getattr(shape, "rotation", 0) or 0),
    }


def _convert_text_shape(shape, theme: ThemeColors) -> dict[str, Any] | None:
    try:
        if not shape.has_text_frame:
            return None
    except Exception:
        return None
    tf = shape.text_frame
    text = (tf.text or "").strip()
    is_placeholder = bool(getattr(shape, "is_placeholder", False))
    if not text and not is_placeholder:
        return None

    frame = _frame_from_shape(shape)
    role = _infer_text_role(tf, frame)
    style = _font_from_text_frame(tf, theme.text)

    placeholders = {
        "title": "Заголовок слайда",
        "subtitle": "Подзаголовок",
        "bulletList": "• Тезис 1\n• Тезис 2\n• Тезис 3",
        "body": "Текст слайда",
    }
    behavior_kind = "placeholder" if (is_placeholder or not text) else "static"
    el = {
        "id": _new_id("el"),
        "type": "text",
        "role": role,
        "contentBehavior": {
            "kind": behavior_kind,
            "readonly": False,
            "fillRole": role,
        },
        "text": "" if behavior_kind == "placeholder" else text,
        "placeholder": placeholders.get(role, "Текст") if behavior_kind == "placeholder" else None,
        "frame": frame,
        "style": style,
        "zIndex": 10,
        "locked": False,
        "visible": True,
    }
    return el


def _convert_picture_shape(shape) -> dict[str, Any]:
    return {
        "id": _new_id("el"),
        "type": "image",
        "role": "image",
        "contentBehavior": {
            "kind": "placeholder",
            "readonly": False,
            "fillRole": "image",
        },
        "assetId": None,
        "fit": "cover",
        "placeholder": "Изображение слайда",
        "frame": _frame_from_shape(shape),
        "zIndex": 5,
        "locked": False,
        "visible": True,
    }


def _convert_auto_shape(shape, theme: ThemeColors) -> dict[str, Any] | None:
    try:
        from pptx.enum.shapes import MSO_SHAPE
        kind_map = {
            MSO_SHAPE.RECTANGLE: "rect",
            MSO_SHAPE.ROUNDED_RECTANGLE: "roundRect",
            MSO_SHAPE.OVAL: "ellipse",
            MSO_SHAPE.ISOSCELES_TRIANGLE: "triangle",
        }
        auto = getattr(shape, "auto_shape_type", None)
        kind = kind_map.get(auto, "rect")
    except Exception:
        kind = "rect"

    fill_hex = None
    stroke_hex = None
    try:
        if shape.fill.type == 1:
            fill_hex = _hex_from_rgb(shape.fill.fore_color.rgb)
    except Exception:
        pass
    try:
        if shape.line and shape.line.color.type:
            stroke_hex = _hex_from_rgb(shape.line.color.rgb)
    except Exception:
        pass

    return {
        "id": _new_id("el"),
        "type": "shape",
        "shape": kind,
        "contentBehavior": {"kind": "static", "readonly": False},
        "frame": _frame_from_shape(shape),
        "style": {
            "fill": fill_hex or theme.primary,
            "stroke": stroke_hex,
            "strokeWidth": 12700 if stroke_hex else None,
        },
        "zIndex": 1,
        "locked": False,
        "visible": True,
    }


def _convert_connector(shape, theme: ThemeColors) -> dict[str, Any]:
    color = None
    width = 12700
    try:
        if shape.line and shape.line.color.type:
            color = _hex_from_rgb(shape.line.color.rgb)
        if shape.line.width:
            width = int(shape.line.width)
    except Exception:
        pass
    return {
        "id": _new_id("el"),
        "type": "line",
        "contentBehavior": {"kind": "static", "readonly": False},
        "frame": _frame_from_shape(shape),
        "style": {
            "color": color or theme.text,
            "widthEmu": max(6350, width),
        },
        "zIndex": 1,
        "locked": False,
        "visible": True,
    }


def _convert_shape(shape, theme: ThemeColors) -> dict[str, Any] | None:
    try:
        shape_type = shape.shape_type
    except Exception:
        shape_type = None

    try:
        if shape.has_text_frame and (shape.text_frame.text or "").strip():
            return _convert_text_shape(shape, theme)
    except Exception:
        pass

    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return _convert_picture_shape(shape)
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return _convert_auto_shape(shape, theme)
    if shape_type == MSO_SHAPE_TYPE.LINE:
        return _convert_connector(shape, theme)
    if shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return _convert_text_shape(shape, theme)
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        results: list[dict[str, Any]] = []
        for sub in shape.shapes:
            converted = _convert_shape(sub, theme)
            if converted:
                results.append(converted)
        return {"__group__": results} if results else None

    try:
        if shape.has_text_frame:
            return _convert_text_shape(shape, theme)
    except Exception:
        pass
    return None


def _flatten_groups(elements: list[dict[str, Any]]) -> list[dict[str, Any]]:
    flat: list[dict[str, Any]] = []
    for el in elements:
        if "__group__" in el:
            flat.extend(_flatten_groups(el["__group__"]))
        else:
            flat.append(el)
    return flat


def _classify_layout(elements: list[dict[str, Any]], slide_index: int, total: int) -> str:
    has_title = any(e.get("role") == "title" for e in elements if e.get("type") == "text")
    has_bullets = any(e.get("role") == "bulletList" for e in elements if e.get("type") == "text")
    has_image = any(e.get("type") == "image" for e in elements)
    text_only_one_block = (
        sum(1 for e in elements if e.get("type") == "text") == 1
    )
    if slide_index == 0:
        return "title"
    if slide_index == total - 1:
        return "conclusion"
    if has_image:
        return "image"
    if has_bullets:
        return "bullets"
    if has_title and text_only_one_block:
        return "title"
    return "text"


def import_pptx_as_template(
    file_bytes: bytes, name: str | None = None
) -> TemplateDocument:
    prs = Presentation(io.BytesIO(file_bytes))
    width = int(prs.slide_width or 12192000)
    height = int(prs.slide_height or 6858000)
    ratio = "16:9" if abs(width / max(1, height) - 16 / 9) < 0.05 else (
        "4:3" if abs(width / max(1, height) - 4 / 3) < 0.05 else "custom"
    )
    colors = _theme_colors_from_pptx(prs)
    theme = PresentationTheme(
        fonts=ThemeFonts(heading="Inter", body="Inter"),
        colors=colors,
    )

    layouts: list[LayoutDocument] = []
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides):
        bg = _slide_background_color(slide, colors.background)
        elements_raw: list[dict[str, Any]] = []
        for shape in slide.shapes:
            converted = _convert_shape(shape, colors)
            if converted:
                elements_raw.append(converted)
        elements = _flatten_groups(elements_raw)
        slide_type = _classify_layout(elements, idx, total)
        layout = LayoutDocument.model_validate(
            {
                "id": _new_id("layout"),
                "name": f"Слайд {idx + 1}",
                "slideType": slide_type,
                "background": bg.model_dump(mode="json"),
                "elements": elements,
            }
        )
        layouts.append(layout)

    if not layouts:
        # Fallback: at least one empty title layout
        layouts.append(
            LayoutDocument.model_validate(
                {
                    "id": _new_id("layout"),
                    "name": "Title",
                    "slideType": "title",
                    "background": {"type": "color", "value": colors.background},
                    "elements": [],
                }
            )
        )

    template = TemplateDocument(
        schemaVersion="1.0.0",
        documentType="template",
        id=_new_id("tmpl"),
        name=(name or "Импорт PPTX").strip()[:80] or "Импорт PPTX",
        slideSize=SlideSize(widthEmu=width, heightEmu=height, ratio=ratio),
        theme=theme,
        layouts=layouts,
        assets=[],
    )
    return template
